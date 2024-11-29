import random
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# List of random font families for headers and body
HEADER_FONTS = ["Arial Black", "Georgia", "Impact", "Verdana", "Tahoma"]
BODY_FONTS = ["Calibri", "Times New Roman", "Courier New", "Garamond", "Franklin Gothic Medium"]

def get_placeholder_font_properties(shape):
    """Retrieve font size and font color from the first run in the shape."""
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            font_size = run.font.size.pt if run.font.size else None
            font_color = run.font.color.rgb if run.font.color and run.font.color.rgb else None
            return font_size, font_color
    return None, None

def set_font_properties(shape, font_size, font_color, font_family):
    """Apply font size, font color, and font family to the text runs."""
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size) if font_size else run.font.size
            if font_color:
                run.font.color.rgb = font_color
            run.font.name = font_family  # Set the font family

def replace_text_in_slide(slide, marker_text, new_text, font_family):
    """Replace text in a slide and keep the placeholder font properties."""
    for shape in slide.shapes:
        if shape.has_text_frame and marker_text in shape.text:
            # Retrieve the placeholder font properties
            font_size, font_color = get_placeholder_font_properties(shape)
            shape.text = new_text  # Replace the text
            set_font_properties(shape, font_size, font_color, font_family)

def GenPPT(input_file, output_file, details, slideInfo):
    # Randomly choose a header and body font family
    header_font_family = random.choice(HEADER_FONTS)
    body_font_family = random.choice(BODY_FONTS)

    # Load the existing PPT (use the path of your saved template)
    presentation = Presentation(input_file)

    # Slide 1: Problem Statement
    problemStatementIndex = slideInfo["p"]
    replace_text_in_slide(presentation.slides[problemStatementIndex], "@@PS1T$$", details["ps1t"], header_font_family)
    replace_text_in_slide(presentation.slides[problemStatementIndex], "@@PS1$$", details["ps1"], body_font_family)
    replace_text_in_slide(presentation.slides[problemStatementIndex], "@@PS2T$$", details["ps2t"], header_font_family)
    replace_text_in_slide(presentation.slides[problemStatementIndex], "@@PS2$$", details["ps2"], body_font_family)

    # Slide 2: Solution
    solutionIndex = slideInfo["s"]
    replace_text_in_slide(presentation.slides[solutionIndex], "@@S1t$$", details["s1t"], header_font_family)
    replace_text_in_slide(presentation.slides[solutionIndex], "@@S1$$", details["s1"], body_font_family)
    replace_text_in_slide(presentation.slides[solutionIndex], "@@S2t$$", details["s2t"], header_font_family)
    replace_text_in_slide(presentation.slides[solutionIndex], "@@S2$$", details["s2"], body_font_family)

    # Slide 3: Market Opportunity
    marketIndex = slideInfo["m"]
    replace_text_in_slide(presentation.slides[marketIndex], "@@MO1t$$", details["m1t"], header_font_family)
    replace_text_in_slide(presentation.slides[marketIndex], "@@MO1$$", details["m1"], body_font_family)
    replace_text_in_slide(presentation.slides[marketIndex], "@@MO2t$$", details["m2t"], header_font_family)
    replace_text_in_slide(presentation.slides[marketIndex], "@@MO2$$", details["m2"], body_font_family)
    replace_text_in_slide(presentation.slides[marketIndex], "@@MO3t$$", details["m3t"], header_font_family)
    replace_text_in_slide(presentation.slides[marketIndex], "@@MO3$$", details["m3"], body_font_family)
    replace_text_in_slide(presentation.slides[marketIndex], "@@MO4t$$", details["m4t"], header_font_family)
    replace_text_in_slide(presentation.slides[marketIndex], "@@MO4$$", details["m4"], body_font_family)

    # Slide 4: Business Model
    businessIndex = slideInfo["b"]
    replace_text_in_slide(presentation.slides[businessIndex], "@@B1$$", details["bm"], body_font_family)

    # Save the modified presentation
    presentation.save(output_file)
    return True
